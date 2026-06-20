"""Deep-dive a client PPTX: slide size, layouts, per-slide shape anatomy.

Feeds deck.py design: where the text boxes sit, what fonts/sizes, what design
elements (images/shapes) repeat per slide, what lives on layout vs slide.
"""

import sys
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

sys.path.insert(0, str(Path(__file__).parent))
from krutidev import krutidev_to_unicode  # noqa: E402

RES = Path(__file__).resolve().parents[1] / "resources"


def cm(v):
    return round(Emu(v).cm, 2) if v is not None else None


def shape_text(shape, limit=70):
    if not shape.has_text_frame:
        return ""
    parts = []
    for para in shape.text_frame.paragraphs:
        seg = ""
        for run in para.runs:
            t = run.text
            if run.font.name and "kruti" in run.font.name.lower():
                t = krutidev_to_unicode(t)
            seg += t
        parts.append(seg)
    return " | ".join(p for p in parts if p)[:limit]


def describe_slide(slide, idx):
    print(f"\n--- slide {idx} (layout: {slide.slide_layout.name!r}) ---")
    for sh in slide.shapes:
        kind = sh.shape_type
        info = f"  [{kind}] '{sh.name}' pos=({cm(sh.left)},{cm(sh.top)})cm size=({cm(sh.width)}x{cm(sh.height)})cm"
        if sh.has_text_frame:
            sizes = Counter()
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.text.strip():
                        sizes[(r.font.name, r.font.size.pt if r.font.size else None,
                               r.font.bold,
                               str(r.font.color.rgb) if r.font.color and r.font.color.type is not None else None)] += 1
            info += f"\n      fonts: {dict(sizes.most_common(4))}"
            info += f"\n      text: {shape_text(sh)!r}"
        print(info)


def main():
    sys.stdout.reconfigure(encoding="utf-8")
    path = RES / (sys.argv[1] if len(sys.argv) > 1 else "ppt target revision series 02.pptx")
    prs = Presentation(str(path))
    print(f"slide size: {cm(prs.slide_width)} x {cm(prs.slide_height)} cm")
    print(f"slides: {len(prs.slides)}")
    layout_use = Counter(s.slide_layout.name for s in prs.slides)
    print(f"layouts used: {dict(layout_use)}")
    for li, layout in enumerate(prs.slide_masters[0].slide_layouts):
        n_img = sum(1 for sh in layout.shapes if sh.shape_type == 13)
        print(f"  layout {li}: {layout.name!r} shapes={len(layout.shapes)} pictures={n_img}")
    master = prs.slide_masters[0]
    print(f"master shapes: {len(master.shapes)}, pictures on master: {sum(1 for s in master.shapes if s.shape_type == 13)}")

    for idx in (1, 2, 3):
        describe_slide(prs.slides[idx - 1], idx)


if __name__ == "__main__":
    main()
