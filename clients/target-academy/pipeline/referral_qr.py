"""Referral QR block for papers and decks.

Every question paper and slide deck the institute hands out becomes a
self-registration point: a small "get your referral link" block with a QR of
the institute's kit page. Students scan, enter their phone, and get their own
share card, no work for the owner. Distribution rides on documents students
already receive.

The kit URL comes from the REFERRAL_KIT_URL env var so it is per-institute and
never hardcoded in the layout code. The default is Target Academy's live kit
link (the pilot), so a build with no env set still works.

Only qrcode + Pillow are used (both in the worker image). This module renders
the QR PNG to a temp file and exposes helpers for python-docx and python-pptx.
"""

import os
import tempfile
from pathlib import Path

import qrcode

# The kit URL is per-institute and secret (it carries the institute's kit
# token), so it MUST come from the REFERRAL_KIT_URL env var and is never
# committed here. If it is unset or malformed, kit_url() returns None and the
# QR block is skipped (the builders treat that as best-effort, no crash).

# English only (product rule): the caption on the block. The paper body is
# Hindi/Kruti-Dev, but this referral block is DSideOS product surface.
CAPTION = "Refer a friend and earn a reward"
SUBCAPTION = "Scan to get your personal referral link"


def kit_url() -> str | None:
    """The configured kit URL, or None if unset/invalid (block is then skipped)."""
    raw = os.environ.get("REFERRAL_KIT_URL", "").strip()
    if not raw:
        return None
    # Must be a well-formed https kit link, or the QR would be a dead link.
    from urllib.parse import urlparse
    try:
        u = urlparse(raw)
    except ValueError:
        return None
    if u.scheme != "https" or not u.netloc or "/kit/" not in u.path:
        return None
    return raw


def render_qr_png(url: str | None = None, box_size: int = 10, border: int = 2) -> Path:
    """Render the kit-URL QR to a temp PNG and return its path.

    High error correction so the code still scans if a printed page is smudged
    or the QR is slightly small on paper. Raises ValueError if no kit URL is
    configured, so the caller skips the block instead of shipping a dead QR.
    """
    target = url or kit_url()
    if not target:
        raise ValueError("no REFERRAL_KIT_URL configured; skipping referral QR")
    qr = qrcode.QRCode(
        version=None,
        error_correction=qrcode.constants.ERROR_CORRECT_H,
        box_size=box_size,
        border=border,
    )
    qr.add_data(target)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")
    fd, name = tempfile.mkstemp(suffix=".png", prefix="referral_qr_")
    os.close(fd)
    out = Path(name)
    # If the save raises, remove the empty temp file we just created so it does
    # not leak (the caller's finally cannot clean up a path it never received).
    try:
        img.save(out)
    except Exception:
        out.unlink(missing_ok=True)
        raise
    return out


# ── python-docx: append a referral block on the last page ──────────────────

def add_referral_block_docx(doc, qr_width_cm: float = 3.0) -> None:
    """Append a centered referral block (caption + QR) after the document's
    existing content. Adds a page break so it sits cleanly on its own, never
    disturbing the question layout, borders or watermark above it."""
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.shared import Cm, Pt, RGBColor

    qr_png = render_qr_png()
    try:
        # Page break so the block starts fresh, not jammed under the answer key.
        brk = doc.add_paragraph()
        brk_ppr = brk._p.get_or_add_pPr()
        brk_ppr.insert(0, brk_ppr.makeelement(qn("w:pageBreakBefore"), {}))

        cap = doc.add_paragraph()
        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = cap.add_run(CAPTION)
        run.bold = True
        run.font.size = Pt(15)
        run.font.name = "Calibri"
        run.font.color.rgb = RGBColor(0x1C, 0x18, 0x30)

        sub = doc.add_paragraph()
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        srun = sub.add_run(SUBCAPTION)
        srun.font.size = Pt(11)
        srun.font.name = "Calibri"
        srun.font.color.rgb = RGBColor(0x56, 0x51, 0x75)

        pic = doc.add_paragraph()
        pic.alignment = WD_ALIGN_PARAGRAPH.CENTER
        pic.add_run().add_picture(str(qr_png), width=Cm(qr_width_cm))
    finally:
        qr_png.unlink(missing_ok=True)


# ── python-pptx: add a final referral slide ────────────────────────────────

def add_referral_slide_pptx(prs, blank_layout=None) -> None:
    """Add a clean final slide with the caption and QR, centered."""
    from pptx.util import Emu, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    layout = blank_layout or prs.slide_masters[0].slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    sw, sh = prs.slide_width, prs.slide_height

    qr_png = render_qr_png()
    try:
        qr_side = Emu(int(sh * 0.42))
        qr_left = Emu(int((sw - qr_side) / 2))
        qr_top = Emu(int(sh * 0.30))
        slide.shapes.add_picture(str(qr_png), qr_left, qr_top, height=qr_side)

        # Caption above the QR.
        box = slide.shapes.add_textbox(Emu(0), Emu(int(sh * 0.14)), sw, Emu(int(sh * 0.14)))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = CAPTION
        r.font.size = Pt(30)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Sub-caption below the QR.
        sub_top = Emu(int(sh * 0.30 + sh * 0.42 + sh * 0.02))
        sbox = slide.shapes.add_textbox(Emu(0), sub_top, sw, Emu(int(sh * 0.10)))
        stf = sbox.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        sr = sp.add_run()
        sr.text = SUBCAPTION
        sr.font.size = Pt(18)
        sr.font.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
    finally:
        qr_png.unlink(missing_ok=True)
