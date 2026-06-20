"""Build a question PPT in the institute's deck template from structured questions.

Template anatomy (from 'ppt target revision series 02.pptx'):
16:9 (33.87 x 19.05 cm), black background (theme bg1->dk1), one question per
slide: full-width banner image strip on top (2.21 cm tall) + one big textbox,
Kruti Dev 010 at 28-54 pt (sized to content), white text with gold (FFC000)
question stem; Calibri for Latin fragments.

Approach: open the client's own deck (inherits master/theme/background), grab
the banner image + geometry from slide 1, delete all original slides, then add
one slide per question on the Blank layout (same master = same black bg).

Usage: python build_deck.py [questions.json] [out.pptx]
"""

import io
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt

sys.path.insert(0, str(Path(__file__).parent))
from krutidev import krutidev_to_unicode, unicode_to_krutidev_runs  # noqa: E402

BASE = Path(__file__).resolve().parents[1]
TEMPLATE_PPTX = BASE / "resources" / "ppt target revision series 02.pptx"

KRUTI = "Kruti Dev 010"
LATIN = "Calibri"
GOLD = RGBColor(0xFF, 0xC0, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W_CM, SLIDE_H_CM = 33.87, 19.05
BANNER_H_CM = 2.21
BOX_TOP_CM = 2.4
BOX_SIDE_INSET_CM = 0.8
EMU_PER_CM = 360000

# image-resolution base: set per build() call so JSON "image" paths resolve
# relative to the questions file's own folder (crops live beside it)
_IMG_BASE = BASE


def _img_path(rel):
    """Resolve a JSON image reference against the questions-file folder."""
    p = Path(rel)
    return p if p.is_absolute() else (_IMG_BASE / p)


def _png_size_px(path):
    from PIL import Image
    with Image.open(path) as im:
        return im.size  # (w, h)


def add_mixed(para, text, size, color, bold=False):
    for run_text, is_latin in unicode_to_krutidev_runs(text):
        r = para.add_run()
        r.text = run_text
        r.font.name = LATIN if is_latin else KRUTI
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color


def question_char_load(q) -> int:
    n = len(q["stem"])
    n += sum(len(s) + 4 for s in q.get("statements") or [])
    n += sum(len(a) + len(b) + 6 for a, b in q.get("match") or [])
    n += len(q.get("lead_in") or "")
    n += sum(len(o) + 5 for o in (q.get("options") or []))
    return n


def font_size_for(load: int) -> int:
    if load <= 160:
        return 40
    if load <= 260:
        return 34
    if load <= 380:
        return 28
    return 24


def grab_banner(prs):
    """Banner image bytes + geometry from the first original slide."""
    for sh in prs.slides[0].shapes:
        if sh.shape_type == 13:  # PICTURE
            return sh.image.blob, sh.left, sh.top, sh.width, sh.height
    return None


def delete_all_slides(prs):
    sld_ids = prs.slides._sldIdLst
    for sld in list(sld_ids):
        prs.part.drop_rel(sld.get(qn("r:id")))
        sld_ids.remove(sld)


def blank_layout(prs):
    for layout in prs.slide_masters[0].slide_layouts:
        if layout.name == "Blank":
            return layout
    return prs.slide_masters[0].slide_layouts[6]


def new_slide(prs, layout, banner):
    slide = prs.slides.add_slide(layout)
    # remove any placeholders the layout dragged in
    for sh in list(slide.shapes):
        sh._element.getparent().remove(sh._element)
    if banner:
        blob, left, top, width, height = banner
        slide.shapes.add_picture(io.BytesIO(blob), left, top, width, height)
    return slide


def _place_image(slide, path, left_cm, top_cm, max_w_cm, max_h_cm):
    """Add a picture scaled to fit (max_w_cm x max_h_cm), centered in that box.
    Returns the bottom edge (cm) so callers can stack below it."""
    pw, ph = _png_size_px(path)
    aspect = pw / ph
    w_cm, h_cm = max_w_cm, max_w_cm / aspect
    if h_cm > max_h_cm:
        h_cm, w_cm = max_h_cm, max_h_cm * aspect
    x_cm = left_cm + (max_w_cm - w_cm) / 2
    slide.shapes.add_picture(str(path), Cm(x_cm), Cm(top_cm), Cm(w_cm), Cm(h_cm))
    return top_cm + h_cm


def add_question_slide(prs, layout, banner, q):
    slide = new_slide(prs, layout, banner)
    has_fig = bool(q.get("image"))
    has_opt_imgs = bool(q.get("option_images"))
    full_w = SLIDE_W_CM - 2 * BOX_SIDE_INSET_CM

    # stem textbox (auto-height; we just give it the top strip)
    box = slide.shapes.add_textbox(
        Cm(BOX_SIDE_INSET_CM), Cm(BOX_TOP_CM),
        Cm(full_w), Cm(SLIDE_H_CM - BOX_TOP_CM - 0.4),
    )
    tf = box.text_frame
    tf.word_wrap = True
    size = font_size_for(question_char_load(q))

    para = tf.paragraphs[0]
    para.space_after = Pt(8)
    add_mixed(para, f"{q['n']}. {q['stem']}", size, GOLD, bold=True)

    for i, st in enumerate(q.get("statements") or [], start=1):
        p = tf.add_paragraph()
        p.space_after = Pt(4)
        add_mixed(p, f"{i}. {st}", size, WHITE)
    for left_t, right_t in q.get("match") or []:
        p = tf.add_paragraph()
        p.space_after = Pt(4)
        add_mixed(p, f"{left_t}  -  {right_t}", size, WHITE)
    if q.get("lead_in"):
        p = tf.add_paragraph()
        p.space_after = Pt(6)
        add_mixed(p, q["lead_in"], size, WHITE)

    # vertical cursor (cm) for anything placed below the stem text
    n_lines = 1 + len(q.get("statements") or []) + len(q.get("match") or []) + \
        (1 if q.get("lead_in") else 0)
    cursor = BOX_TOP_CM + n_lines * (size / 28) * 1.1 + 0.6

    # the question's own figure (Q3-type: one diagram under the stem)
    if has_fig:
        fig_max_h = (SLIDE_H_CM - cursor - 0.5) * (0.55 if has_opt_imgs else 0.75)
        cursor = _place_image(slide, _img_path(q["image"]),
                              BOX_SIDE_INSET_CM, cursor, full_w, fig_max_h) + 0.4

    labels = ["(a)", "(b)", "(c)", "(d)", "(e)"]

    # option diagrams (Q60-type: each (a)-(d) is a cropped figure, laid in a row)
    if has_opt_imgs:
        imgs = q["option_images"]
        cols = len(imgs)
        gap = 0.4
        cell_w = (full_w - gap * (cols - 1)) / cols
        cell_h = SLIDE_H_CM - cursor - 0.4
        for i, rel in enumerate(imgs):
            x = BOX_SIDE_INSET_CM + i * (cell_w + gap)
            lbl = slide.shapes.add_textbox(Cm(x), Cm(cursor), Cm(cell_w), Cm(0.9))
            lp = lbl.text_frame.paragraphs[0]
            add_mixed(lp, labels[i], 24, GOLD, bold=True)
            _place_image(slide, _img_path(rel), x, cursor + 0.9, cell_w, cell_h - 0.9)
        return

    # text options (normal path)
    opts = q["options"]
    if not q.get("long_options") and len(opts) == 4 and all(len(o) <= 20 for o in opts):
        for i in (0, 2):
            p = tf.add_paragraph()
            p.space_after = Pt(4)
            add_mixed(p, f"{labels[i]} {opts[i]}\t\t{labels[i + 1]} {opts[i + 1]}", size, WHITE)
    else:
        for i, o in enumerate(opts):
            p = tf.add_paragraph()
            p.space_after = Pt(4)
            add_mixed(p, f"{labels[i]} {o}", size, WHITE)


def add_answer_slide(prs, layout, banner, questions):
    slide = new_slide(prs, layout, banner)
    box = slide.shapes.add_textbox(
        Cm(BOX_SIDE_INSET_CM), Cm(BOX_TOP_CM),
        Cm(SLIDE_W_CM - 2 * BOX_SIDE_INSET_CM), Cm(SLIDE_H_CM - BOX_TOP_CM - 0.4),
    )
    tf = box.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.space_after = Pt(12)
    add_mixed(para, "उत्तरमाला", 40, GOLD, bold=True)
    row = []
    for q in questions:
        row.append(f"{q['n']}. ({q['answer']})")
        if len(row) == 5:
            p = tf.add_paragraph()
            p.space_after = Pt(6)
            add_mixed(p, "      ".join(row), 28, WHITE)
            row = []
    if row:
        p = tf.add_paragraph()
        add_mixed(p, "      ".join(row), 28, WHITE)


def build(questions_path: Path, out_path: Path, answer_key: bool = True):
    """answer_key=False -> verbatim-MCQ deck shown live in class (no answers).
    answer_key=True  -> practice deck with the trailing उत्तरमाला slide."""
    global _IMG_BASE
    _IMG_BASE = questions_path.resolve().parent  # crops resolve beside the JSON
    data = json.loads(questions_path.read_text(encoding="utf-8"))
    prs = Presentation(str(TEMPLATE_PPTX))
    banner = grab_banner(prs)
    layout = blank_layout(prs)
    delete_all_slides(prs)
    for q in data["questions"]:
        add_question_slide(prs, layout, banner, q)
    if answer_key:
        add_answer_slide(prs, layout, banner, data["questions"])
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"OK: {out_path}")


def verify(out_path: Path):
    prs = Presentation(str(out_path))
    n = len(prs.slides)
    sample = ""
    for sh in prs.slides[0].shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    t = r.text
                    if r.font.name and "kruti" in (r.font.name or "").lower():
                        t = krutidev_to_unicode(t)
                    sample += t
            break
    pics = sum(1 for s in prs.slides for sh in s.shapes if sh.shape_type == 13)
    print(f"verify: {n} slides, {pics} banner pictures, sample: {sample[:70]!r}")


if __name__ == "__main__":
    sys.stdout.reconfigure(encoding="utf-8")
    q = Path(sys.argv[1]) if len(sys.argv) > 1 else BASE / "review" / "spike3-questions.json"
    data = json.loads(q.read_text(encoding="utf-8"))
    default_name = data.get("ppt_filename") or (
        Path(data.get("filename", "Practice Paper - Target Academy.docx")).stem + ".pptx"
    )
    out = Path(sys.argv[2]) if len(sys.argv) > 2 else BASE / "review" / default_name
    # pass "no-key" as a 3rd arg for the in-class verbatim-MCQ deck
    answer_key = "no-key" not in sys.argv[3:]
    build(q, out, answer_key=answer_key)
    verify(out)
