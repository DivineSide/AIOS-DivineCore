"""Inspect the client's raw materials: formats, fonts, encoding (Kruti vs Unicode).

First-look tool for Phase 2 ingest — reports per file: fonts used, sample text
(raw + KrutiDev->Unicode converted when the text looks Kruti-encoded), so we can
confirm the encoding assumption before building ingest.py proper.
"""

import sys
from collections import Counter
from pathlib import Path

from docx import Document
from pptx import Presentation

from krutidev import krutidev_to_unicode

RES = Path(__file__).resolve().parents[1] / "resources"

DEVANAGARI = range(0x0900, 0x0980)


def classify(text: str) -> str:
    """Heuristic: is this run Unicode Devanagari, probable-Kruti ASCII, or Latin?"""
    if any(ord(c) in DEVANAGARI for c in text):
        return "unicode-devanagari"
    # Kruti tell-tales: heavy use of these glyph codepoints in "English" text
    telltales = sum(text.count(c) for c in "fdkjghqwlu;~`¼½ÞßAÙÆÇÉÊ{}=+")
    letters = sum(c.isalpha() for c in text)
    if letters and telltales / max(letters, 1) > 0.15:
        return "probable-krutidev"
    return "latin/other"


def sample_runs_pptx(path: Path, limit=12):
    prs = Presentation(path)
    fonts, samples = Counter(), []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    t = run.text.strip()
                    if not t:
                        continue
                    fonts[run.font.name or "(inherit)"] += 1
                    if len(samples) < limit and len(t) > 3:
                        samples.append((f"slide{i+1}", run.font.name, t[:60]))
    return len(prs.slides), fonts, samples


def sample_runs_docx(path: Path, limit=12):
    doc = Document(str(path))
    fonts, samples = Counter(), []
    paras = list(doc.paragraphs)
    for tbl in doc.tables:
        for row in tbl.rows:
            for cell in row.cells:
                paras.extend(cell.paragraphs)
    for para in paras:
        for run in para.runs:
            t = run.text.strip()
            if not t:
                continue
            fonts[run.font.name or "(inherit)"] += 1
            if len(samples) < limit and len(t) > 3:
                samples.append(("", run.font.name, t[:60]))
    return len(paras), fonts, samples


def report(path: Path):
    print(f"\n{'='*70}\n{path.name}\n{'='*70}")
    if path.suffix.lower() == ".pptx":
        n, fonts, samples = sample_runs_pptx(path)
        print(f"slides: {n}")
    elif path.suffix.lower() == ".docx":
        n, fonts, samples = sample_runs_docx(path)
        print(f"paragraphs (incl. tables): {n}")
    else:
        print("(image/other — skipped)")
        return
    print(f"fonts: {dict(fonts.most_common(8))}")
    for where, font, text in samples[:8]:
        kind = classify(text)
        line = f"  [{where} | {font} | {kind}] {text!r}"
        if kind == "probable-krutidev":
            line += f"\n      -> k2u: {krutidev_to_unicode(text)!r}"
        print(line)


if __name__ == "__main__":
    sys.stdout.reconfigure(encoding="utf-8")
    if len(sys.argv) > 1:
        targets = [RES / a for a in sys.argv[1:]]
    else:
        targets = sorted(RES.glob("*"))
    for p in targets:
        if p.is_file():
            report(p)
