"""Spike: Devanagari rendering via python-pptx (go/no-go check, brain/spec.md edge cases).

Generates review/spike_devanagari.pptx with Hindi text stressing the known sharp edges:
conjuncts, matras, danda, mixed Hindi-English lines, bold/size runs, and two candidate
fonts (Nirmala UI, Mangal). Pass = everything renders correctly when opened in PowerPoint.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parents[1] / "review" / "spike_devanagari.pptx"

# Stress content: conjuncts (क्ष त्र ज्ञ द्ध श्र), matras, danda, numerals, mixed-script
TITLE = "भारतीय संविधान — मौलिक अधिकार"
BULLETS = [
    "अनुच्छेद 14: विधि के समक्ष समता का अधिकार।",
    "क्षत्रिय, ज्ञान, द्धन्द्व, श्रद्धा — संयुक्ताक्षर परीक्षण।",
    "मिश्रित पंक्ति: Fundamental Rights (मौलिक अधिकार) — Article 19(1)(a)।",
    "प्रश्न-पत्र खण्ड 'अ' — प्रत्येक प्रश्न 2 अंक का है।",
    "हिन्दी अंक: १२३४५६७८९० तथा मानक अंक: 1234567890",
]

prs = Presentation()

for font_name in ("Nirmala UI", "Mangal"):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
    slide.shapes.title.text = f"{TITLE}  [{font_name}]"
    for run in slide.shapes.title.text_frame.paragraphs[0].runs:
        run.font.name = font_name
        run.font.size = Pt(32)

    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(BULLETS):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(18)
        run.font.bold = i == 1  # bold run on the conjunct line

    # free textbox (deck.py will also place text outside placeholders)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(6.4), Inches(9), Inches(0.6))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = "टेक्स्टबॉक्स परीक्षण — यह पंक्ति प्लेसहोल्डर के बाहर है।"
    r.font.name = font_name
    r.font.size = Pt(14)

# Slide 3 — Kruti Dev 010 legacy-encoding check. Text below is ASCII in Kruti Dev
# encoding (NOT Unicode): "Hkkjr" → भारत, "fgUnh" → हिन्दी. Pass = renders as Hindi.
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Kruti Dev 010 encoding test"
body = slide.placeholders[1].text_frame
body.clear()
for i, (legacy, expected) in enumerate([("Hkkjr", "भारत"), ("fgUnh", "हिन्दी")]):
    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
    run = p.add_run()
    run.text = legacy
    run.font.name = "Kruti Dev 010"
    run.font.size = Pt(40)
    note = p.add_run()
    note.text = f"   <- should read {expected}"
    note.font.name = "Calibri"
    note.font.size = Pt(14)

# Slide 4 — full pipeline path: Unicode (what the LLM generates) -> krutidev.py
# converter -> Kruti Dev 010 font runs, with Latin words in Calibri.
from krutidev import unicode_to_krutidev_runs  # noqa: E402

UNICODE_LINES = [
    "भारतीय संविधान एवं मौलिक अधिकार।",
    "क्षत्रिय, ज्ञान, श्रद्धा — संयुक्ताक्षर परीक्षण।",
    "मौलिक अधिकार (Fundamental Rights) अनुच्छेद 19 देखें।",
    "प्रश्न-पत्र खण्ड 'अ' में प्रत्येक प्रश्न 2 अंक का है।",
]
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Unicode -> Kruti Dev converter test"
body = slide.placeholders[1].text_frame
body.clear()
for i, line in enumerate(UNICODE_LINES):
    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
    for run_text, is_latin in unicode_to_krutidev_runs(line):
        run = p.add_run()
        run.text = run_text
        run.font.name = "Calibri" if is_latin else "Kruti Dev 010"
        run.font.size = Pt(20)

OUT.parent.mkdir(parents=True, exist_ok=True)
try:
    prs.save(OUT)
except PermissionError:  # file open in PowerPoint — save alongside instead
    OUT = OUT.with_stem(OUT.stem + "_v2")
    prs.save(OUT)
print(f"OK: {OUT}")
