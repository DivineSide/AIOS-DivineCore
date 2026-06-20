"""Ingest the client's raw materials into a Unicode text corpus.

Reads .docx/.pptx from resources/, converts Kruti Dev 010 runs to Unicode
Devanagari, and writes one markdown file per source into corpus/parsed/.

Conversion is done per PARAGRAPH, not per run: DOCX fragments words across
runs mid-glyph, and Kruti positional rules (chhoti-i "f", reph "Z") span those
fragments (see brain/format-conventions.md, ingest lesson #1). Within a
paragraph, contiguous runs of the same script-class are joined, then Kruti
segments are converted and Latin segments passed through.
"""

import sys
from pathlib import Path

from docx import Document
from pptx import Presentation

from krutidev import krutidev_to_unicode

BASE = Path(__file__).resolve().parents[1]
RES = BASE / "resources"
OUT = BASE / "corpus" / "parsed"

DEVANAGARI = range(0x0900, 0x0980)
_KRUTI_TELLTALES = set("fdkjghqwlu;~`¼½ÞßAÙÆÇÉÊ{}=+")


def _is_kruti_font(name) -> bool:
    return bool(name) and "kruti" in name.lower()


def _guess_class(text: str) -> bool:
    """True = treat as Kruti when the font is inherited/unknown."""
    if any(ord(c) in DEVANAGARI for c in text):
        return False  # already Unicode Devanagari — pass through
    letters = sum(c.isalpha() for c in text)
    telltales = sum(c in _KRUTI_TELLTALES for c in text)
    return letters > 0 and telltales / max(letters, 1) > 0.15


def paragraph_text(para) -> str:
    """Join runs into (class, text) segments, convert Kruti segments."""
    segments = []  # list of [is_kruti, text]
    prev_kruti = None
    for run in para.runs:
        t = run.text
        if not t:
            continue
        if run.font.name is not None:
            is_kruti = _is_kruti_font(run.font.name)
        elif prev_kruti is not None:
            is_kruti = prev_kruti  # inherit class from previous run
        else:
            is_kruti = _guess_class(t)
        if segments and segments[-1][0] == is_kruti:
            segments[-1][1] += t
        else:
            segments.append([is_kruti, t])
        prev_kruti = is_kruti
    return "".join(
        krutidev_to_unicode(text) if is_kruti else text for is_kruti, text in segments
    )


def ingest_docx(path: Path) -> str:
    doc = Document(str(path))
    lines = []
    for para in doc.paragraphs:
        t = paragraph_text(para).strip()
        if t:
            lines.append(t)
    for ti, tbl in enumerate(doc.tables):
        lines.append(f"\n<!-- table {ti + 1} -->")
        for row in tbl.rows:
            cells = []
            for cell in row.cells:
                cell_text = " ".join(
                    paragraph_text(p).strip() for p in cell.paragraphs
                ).strip()
                cells.append(cell_text)
            if any(cells):
                lines.append("| " + " | ".join(cells) + " |")
    return "\n".join(lines)


def ingest_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    lines = []
    for i, slide in enumerate(prs.slides):
        lines.append(f"\n## slide {i + 1}")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = paragraph_text(para).strip()
                if t:
                    lines.append(t)
    return "\n".join(lines)


def main():
    sys.stdout.reconfigure(encoding="utf-8")
    OUT.mkdir(parents=True, exist_ok=True)
    for path in sorted(RES.glob("*")):
        if path.suffix.lower() == ".docx":
            body = ingest_docx(path)
        elif path.suffix.lower() == ".pptx":
            body = ingest_pptx(path)
        else:
            continue
        dest = OUT / (path.stem + ".md")
        header = f"<!-- source: {path.name} | ingested via ingest.py (Kruti Dev 010 -> Unicode) -->\n\n"
        dest.write_text(header + body, encoding="utf-8")
        print(f"{path.name} -> {dest.relative_to(BASE)} ({len(body):,} chars)")


if __name__ == "__main__":
    main()
